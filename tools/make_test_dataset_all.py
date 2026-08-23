#!/usr/bin/env python3
"""
make_test_dataset_all.py
========================
Generate a test corpus covering EVERY file type the summarizer can extract,
laid out like a real file share. Reuses the proven builders in
make_test_dataset.py (LAS, SEG-Y, PDF, scout PDF, xlsx, docx, shapefile,
geojson) and adds the rest (KML/KMZ, PowerPoint, CSV/TSV/TXT, ASCII/.prn/.dev
logs, email, OpenDocument, RTF, GeoPackage, MapInfo, GeoTIFF, plain image,
WITSML, OSDU JSON well log).

    python tools/make_test_dataset_all.py --root C:\\Bulk\\TestDataAll
    python tools/make_test_dataset_all.py --root C:\\Bulk\\TestDataAll --validate

--validate runs each generated file back through modules.file_summarizer.summarize()
and prints a coverage table (format · uwi · ok/err), so you can confirm the
pipeline recognizes every type before scanning it for real.

Each generator is wrapped so a missing optional library skips that format
rather than aborting the run. UWIs use the obviously-fake 42999% prefix for
easy cleanup.
"""

import argparse
import os
import sys
import zipfile


# The REPO ROOT, not tools/. Python puts the SCRIPT's own directory on
# sys.path[0], so `python tools/<name>.py` cannot import dataview without
# this. app_v4.py does the same insert; see tools/reconcile_orphans.py.
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import make_test_dataset as base   # noqa: E402  (reuse proven builders)

WELLS = base.WELLS
SURVEYS = base.SURVEYS
api = base.fmt_api


# ── new per-well generators ───────────────────────────────────────────────────
def gen_csv(path, w, sep=","):
    hdr = ["UWI", "WELL_NAME", "API", "OPERATOR", "FIELD", "COUNTY", "STATE"]
    row = [w["uwi"], w["name"], api(w["uwi"]), w["operator"],
           w["field"], w["county"], w["state"]]
    with open(path, "w", newline="") as f:
        f.write(sep.join(hdr) + "\n" + sep.join(str(x) for x in row) + "\n")


def _kml_text(w):
    return (f'<?xml version="1.0"?>\n'
            f'<kml xmlns="http://www.opengis.net/kml/2.2"><Document>\n'
            f'<Placemark><name>{w["name"]}</name>\n'
            f'<ExtendedData>\n'
            f'<Data name="UWI"><value>{w["uwi"]}</value></Data>\n'
            f'<Data name="Operator"><value>{w["operator"]}</value></Data>\n'
            f'<Data name="Field"><value>{w["field"]}</value></Data>\n'
            f'</ExtendedData>\n'
            f'<Point><coordinates>{w["lon"]},{w["lat"]},0</coordinates></Point>\n'
            f'</Placemark></Document></kml>')


def gen_kml(path, w):
    open(path, "w").write(_kml_text(w))


def gen_kmz(path, w):
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("doc.kml", _kml_text(w))


def gen_pptx(path, w):
    from pptx import Presentation
    prs = Presentation()
    sl = prs.slides.add_slide(prs.slide_layouts[1])
    sl.shapes.title.text = f"Well Report: {w['name']}"
    sl.placeholders[1].text = (f"Well Name {w['name']}\nUWI {w['uwi']}\n"
                               f"API {api(w['uwi'])}\nOperator {w['operator']}\n"
                               f"Field {w['field']}")
    prs.save(path)


def gen_ascii_log(path, w):
    lines = [f"# WELL NAME: {w['name']}", f"# UWI: {w['uwi']}",
             f"# API: {api(w['uwi'])}", "DEPT  GR  RHOB  NPHI"]
    top = w["strt"]
    for i in range(30):
        lines.append(f"{top + i:.1f}  {50 + i % 40}  2.35  0.15")
    open(path, "w").write("\n".join(lines) + "\n")


def gen_dev(path, w):
    lines = [f"# WELL {w['name']}", f"# UWI {w['uwi']}", "MD  INC  AZI"]
    for i in range(20):
        lines.append(f"{i * 100}  {min(i * 1.5, 90):.1f}  135.0")
    open(path, "w").write("\n".join(lines) + "\n")


def gen_eml(path, w):
    open(path, "w").write(
        "From: geologist@example.com\nTo: data@example.com\n"
        f"Subject: Log data for UWI {w['uwi']} ({w['name']})\n\n"
        f"Attached are logs for WELL NAME {w['name']}, UWI {w['uwi']}, "
        f"operated by {w['operator']} in the {w['field']} field.\n")


def gen_odf(path, w):
    if path.endswith(".odt"):
        mt, root = "text", "text"
    elif path.endswith(".ods"):
        mt, root = "spreadsheet", "spreadsheet"
    else:
        mt, root = "presentation", "presentation"
    content = (
        "<?xml version='1.0'?><office:document-content "
        "xmlns:office='urn:oasis'><office:body><office:" + root + ">"
        f"<text:p>Well Report</text:p>"
        f"<text:p>Well Name: {w['name']}</text:p>"
        f"<text:p>UWI {w['uwi']} API {api(w['uwi'])}</text:p>"
        f"<text:p>Operator {w['operator']}, Field {w['field']}.</text:p>"
        "</office:" + root + "></office:body></office:document-content>")
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("mimetype",
                   "application/vnd.oasis.opendocument." + mt)
        z.writestr("content.xml", content)


def gen_rtf(path, w):
    open(path, "w").write(
        r"{\rtf1\ansi\deff0 Well Report\par Well Name: " + w["name"] +
        r"\par UWI " + w["uwi"] + r" API " + api(w["uwi"]) +
        r"\par Operator " + w["operator"] + r", Field " + w["field"] + r".\par}")


def gen_witsml(path, w):
    open(path, "w").write(
        '<?xml version="1.0"?>\n'
        '<wells xmlns="http://www.witsml.org/schemas/1series" '
        'version="1.4.1.1">\n'
        f'<well uid="{w["uwi"]}">\n<name>{w["name"]}</name>\n'
        f'<numAPI>{api(w["uwi"])}</numAPI>\n'
        f'<operator>{w["operator"]}</operator>\n<field>{w["field"]}</field>\n'
        f'<county>{w["county"]}</county><state>{w["state"]}</state>\n'
        f'<country>{w["country"]}</country>\n</well>\n</wells>\n')


def gen_json_log(path, w):
    import json
    obj = {
        "kind": "osdu:wks:work-product-component--WellLog:1.0.0",
        "header": {"name": w["name"], "wellbore": w["name"]},
        "data": {"UWI": w["uwi"], "API": api(w["uwi"]), "WellName": w["name"],
                 "Operator": w["operator"], "Field": w["field"],
                 "State": w["state"], "County": w["county"]},
        "curves": [{"name": "GR", "unit": "GAPI"},
                   {"name": "RHOB", "unit": "G/C3"}]}
    open(path, "w").write(json.dumps(obj, indent=2))


# ── new spatial / raster generators ───────────────────────────────────────────
def _well_gdf(wells):
    import geopandas as gpd
    from shapely.geometry import Point
    return gpd.GeoDataFrame(
        {"UWI": [x["uwi"] for x in wells],
         "WELL_NAME": [x["name"] for x in wells],
         "OPERATOR": [x["operator"] for x in wells]},
        geometry=[Point(x["lon"], x["lat"]) for x in wells], crs="EPSG:4326")


def gen_gpkg(path, wells):
    _well_gdf(wells).to_file(path, driver="GPKG")


def gen_mapinfo_tab(stem, wells):
    _well_gdf(wells).to_file(stem + ".tab", driver="MapInfo File")


def gen_mapinfo_mif(stem, wells):
    try:
        _well_gdf(wells).to_file(stem + ".mif", driver="MapInfo File",
                                 FORMAT="MIF")
    except Exception:
        # some GDAL builds need the option spelled differently; fall back to .tab
        _well_gdf(wells).to_file(stem + ".tab", driver="MapInfo File")


def gen_geotiff(path):
    import numpy as np
    import rasterio
    from rasterio.transform import from_bounds
    arr = (np.random.rand(1, 64, 64) * 255).astype("uint8")
    with rasterio.open(path, "w", driver="GTiff", height=64, width=64, count=1,
                       dtype="uint8", crs="EPSG:4326",
                       transform=from_bounds(-96, 29, -95, 30, 64, 64)) as ds:
        ds.write(arr)


def gen_png(path):
    from PIL import Image
    Image.new("RGB", (120, 160), (90, 40, 40)).save(path)


# ── build the corpus ──────────────────────────────────────────────────────────
def build(root):
    results = []

    def attempt(label, fn):
        try:
            fn()
            results.append(("ok", label, None))
        except Exception as e:
            results.append(("skip", label, f"{type(e).__name__}: {e}"))

    wells_dir = os.path.join(root, "Wells")
    for w in WELLS:
        d = os.path.join(wells_dir, f"{w['uwi']}__{w['name'].replace(' ', '_')}")
        os.makedirs(d, exist_ok=True)
        u = w["uwi"]
        P = lambda fn: os.path.join(d, fn)  # noqa: E731
        attempt(f"{u} run1.las", lambda w=w, P=P: base.gen_las(P(f"{w['uwi']}_run1.las"), w, base.RUN1, "run1"))
        attempt(f"{u} run2.las", lambda w=w, P=P: base.gen_las(P(f"{w['uwi']}_run2.las"), w, base.RUN2, "run2"))
        attempt(f"{u} directional.pdf", lambda w=w, P=P: base.gen_directional_pdf(P(f"{w['uwi']}_directional.pdf"), w))
        attempt(f"{u} formation.pdf", lambda w=w, P=P: base.gen_formation_pdf(P(f"{w['uwi']}_formation.pdf"), w))
        attempt(f"{u} scout.pdf", lambda w=w, P=P: base.gen_scout_pdf(P(f"{w['uwi']}_scout.pdf"), w))
        attempt(f"{u} completion.xlsx", lambda w=w, P=P: base.gen_xlsx(P(f"{w['uwi']}_completion.xlsx"), w))
        attempt(f"{u} summary.docx", lambda w=w, P=P: base.gen_docx(P(f"{w['uwi']}_summary.docx"), w))
        attempt(f"{u} deck.pptx", lambda w=w, P=P: gen_pptx(P(f"{w['uwi']}_deck.pptx"), w))
        attempt(f"{u} data.csv", lambda w=w, P=P: gen_csv(P(f"{w['uwi']}_data.csv"), w))
        attempt(f"{u} data.tsv", lambda w=w, P=P: gen_csv(P(f"{w['uwi']}_data.tsv"), w, sep="\t"))
        attempt(f"{u} data.txt", lambda w=w, P=P: gen_csv(P(f"{w['uwi']}_data.txt"), w))
        attempt(f"{u} log.asc", lambda w=w, P=P: gen_ascii_log(P(f"{w['uwi']}_log.asc"), w))
        attempt(f"{u} log.prn", lambda w=w, P=P: gen_ascii_log(P(f"{w['uwi']}_log.prn"), w))
        attempt(f"{u} survey.dev", lambda w=w, P=P: gen_dev(P(f"{w['uwi']}_survey.dev"), w))
        attempt(f"{u} note.eml", lambda w=w, P=P: gen_eml(P(f"{w['uwi']}_note.eml"), w))
        attempt(f"{u} report.odt", lambda w=w, P=P: gen_odf(P(f"{w['uwi']}_report.odt"), w))
        attempt(f"{u} report.ods", lambda w=w, P=P: gen_odf(P(f"{w['uwi']}_report.ods"), w))
        attempt(f"{u} report.odp", lambda w=w, P=P: gen_odf(P(f"{w['uwi']}_report.odp"), w))
        attempt(f"{u} memo.rtf", lambda w=w, P=P: gen_rtf(P(f"{w['uwi']}_memo.rtf"), w))
        attempt(f"{u} location.kml", lambda w=w, P=P: gen_kml(P(f"{w['uwi']}_location.kml"), w))
        attempt(f"{u} location.kmz", lambda w=w, P=P: gen_kmz(P(f"{w['uwi']}_location.kmz"), w))
        attempt(f"{u} well.xml", lambda w=w, P=P: gen_witsml(P(f"{w['uwi']}_well.xml"), w))
        attempt(f"{u} welllog.json", lambda w=w, P=P: gen_json_log(P(f"{w['uwi']}_welllog.json"), w))

    seis_dir = os.path.join(root, "Seismic")
    for s in SURVEYS:
        d = os.path.join(seis_dir, s["file"])
        os.makedirs(d, exist_ok=True)
        attempt(f"{s['name']} {s['ext']}", lambda s=s, d=d: base.gen_segy(os.path.join(d, s["file"] + s["ext"]), s))
        attempt(f"{s['name']} .p190", lambda s=s, d=d: base.gen_p190(os.path.join(d, s["file"] + ".p190"), s))

    sp = os.path.join(root, "Spatial")
    os.makedirs(sp, exist_ok=True)
    attempt("well_locations.shp", lambda: base.gen_shapefile(os.path.join(sp, "well_locations"), WELLS))
    attempt("well_locations.geojson", lambda: base.gen_geojson(os.path.join(sp, "well_locations.geojson"), WELLS))
    attempt("well_locations.gpkg", lambda: gen_gpkg(os.path.join(sp, "well_locations.gpkg"), WELLS))
    attempt("well_locations.tab", lambda: gen_mapinfo_tab(os.path.join(sp, "well_locations_tab"), WELLS))
    attempt("well_locations.mif", lambda: gen_mapinfo_mif(os.path.join(sp, "well_locations_mif"), WELLS))
    attempt("basemap.tif (GeoTIFF)", lambda: gen_geotiff(os.path.join(sp, "basemap.tif")))
    attempt("scan_image.png", lambda: gen_png(os.path.join(sp, "scan_image.png")))

    # manifest
    with open(os.path.join(root, "MANIFEST.txt"), "w") as f:
        f.write("Test corpus — all supported file types\n")
        for status, label, err in results:
            f.write(f"[{status:4}] {label}" + (f"   ({err})" if err else "") + "\n")
    return results


def validate(root):
    try:
        from dataview.file_catalog.file_summarizer import summarize
    except Exception:
        from dataview.file_catalog.file_summarizer import summarize
    rows = []
    for dp, _, files in os.walk(root):
        for fn in files:
            if fn == "MANIFEST.txt":
                continue
            p = os.path.join(dp, fn)
            try:
                s = summarize(p)
                rows.append((os.path.splitext(fn)[1].lower(), s.get("format"),
                             s.get("uwi"), s.get("error")))
            except Exception as e:
                rows.append((os.path.splitext(fn)[1].lower(), "?", None, str(e)))
    return rows


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--root", default="TestDataAll")
    ap.add_argument("--validate", action="store_true",
                    help="run each file back through summarize() and report")
    a = ap.parse_args()

    os.makedirs(a.root, exist_ok=True)
    results = build(a.root)
    ok = sum(1 for r in results if r[0] == "ok")
    sk = sum(1 for r in results if r[0] == "skip")
    print(f"[BUILD] {ok} file(s) generated, {sk} skipped -> {a.root}")
    for status, label, err in results:
        if status == "skip":
            print(f"   SKIP {label}: {err}")

    if a.validate:
        print("\n[VALIDATE] running summarize() on every file:")
        rows = validate(a.root)
        rows.sort()
        got = sum(1 for _, _, u, _ in rows if u)
        for ext, fmt, uwi, err in rows:
            mark = "✓" if uwi else (" " if not err else "✗")
            print(f"  {mark} {ext:9} fmt={str(fmt):12} uwi={str(uwi):16}"
                  + (f" err={err[:40]}" if err else ""))
        print(f"\n[VALIDATE] {got}/{len(rows)} files yielded a UWI "
              "(spatial basemaps/images intentionally have none).")


if __name__ == "__main__":
    main()
