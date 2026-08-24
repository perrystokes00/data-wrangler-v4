"""
modules/file_summarizer.py
==========================
Universal file summarizer for all Data Wrangler supported formats.
Returns a standard summary dict for any file type.

Usage:
    from dataview.file_catalog.file_summarizer import summarize
    info = summarize("/path/to/well.las")
    print(info["description"])
"""
from __future__ import annotations
import re, uuid
from pathlib import Path
from typing import Optional

# ── Standard summary structure ────────────────────────────────────────────────
def _base(file_path: str, fmt: str) -> dict:
    fp = Path(file_path)
    return {
        "file_path":   file_path,
        "file_name":   fp.name,
        "format":      fmt,
        "size_kb":     round(fp.stat().st_size / 1024, 1) if fp.exists() else 0,
        "well_name":   None,
        "uwi":         None,
        "description": "",
        "key_fields":  {},
        "warnings":    [],
        "ppdm_hints":  [],
        "error":       None,
    }


# ══════════════════════════════════════════════════════════════════════════════
# LAS
# ══════════════════════════════════════════════════════════════════════════════
def _las_value(section, *mnems):
    """First non-empty .value among the given LAS mnemonics, else None.

    lasio's section.get() returns a HeaderItem (which is falsy in a boolean
    context and stringifies to its repr), so the value must be read off .value
    explicitly — never via `get(...) or ...` chaining or str(get(...)).
    """
    for m in mnems:
        it = section.get(m)
        v = getattr(it, "value", None)
        if v is None:
            continue
        v = str(v).strip()
        if v and v.lower() not in ("unknown", "none", "--"):
            return v
    return None


def _summarize_las(file_path: str) -> dict:
    s = _base(file_path, "LAS")
    try:
        import lasio
        from dataview.file_catalog.las_reader import read_las
        las = read_las(file_path, ignore_header_errors=True)

        well    = las.well
        curves  = las.curves
        depth   = las.index

        s["well_name"] = _las_value(well, "WELL", "WN", "WELLNAME")
        s["uwi"]       = _las_value(well, "UWI", "API")

        n_curves  = len(curves)
        curve_nms = [c.mnemonic for c in curves
                     if c.mnemonic.upper() not in ("DEPT","DEPTH","MD")][:8]
        d_start   = round(float(depth.min()), 1) if len(depth) else 0
        d_stop    = round(float(depth.max()), 1) if len(depth) else 0
        d_step    = round(float(las.well.get("STEP","").value), 3) if hasattr(
                    las.well.get("STEP",""), "value") else 0
        null_val  = str(las.well.get("NULL","").value) if hasattr(
                    las.well.get("NULL",""), "value") else "-9999.25"
        company   = str(las.well.get("COMP","").value) if hasattr(
                    las.well.get("COMP",""), "value") else ""
        field     = str(las.well.get("FLD","").value) if hasattr(
                    las.well.get("FLD",""), "value") else ""

        s["description"] = (
            f"LAS {las.version.get('VERS','2.0').value if hasattr(las.version.get('VERS',''), 'value') else '2.0'}"
            f" · {n_curves} curves · {d_start:,.0f}–{d_stop:,.0f} ft"
            f" · {d_step} ft step"
            f" · Curves: {', '.join(curve_nms)}"
        )
        s["key_fields"] = {
            "curves":      n_curves,
            "curve_names": curve_nms,
            "depth_start": d_start,
            "depth_stop":  d_stop,
            "depth_step":  d_step,
            "null_value":  null_val,
            "company":     company,
            "field":       field,
            "samples":     len(depth),
        }
        s["ppdm_hints"] = ["dbo.WELL_LOG_SAMPLE", "dbo.WELL"]

        # Warnings
        if d_step == 0:
            s["warnings"].append("Step size is zero or missing")
        if not (s["uwi"] or "").strip():
            s["warnings"].append("No UWI/API found in header")
        if n_curves < 2:
            s["warnings"].append("Only 1 curve — file may be incomplete")

    except Exception as e:
        s["error"] = str(e)
    return s


# ══════════════════════════════════════════════════════════════════════════════
# DLIS
# ══════════════════════════════════════════════════════════════════════════════
def _summarize_dlis(file_path: str) -> dict:
    s = _base(file_path, "DLIS")
    try:
        import dlisio
        with dlisio.dlis(file_path) as (f, *tail):
            lfs = [f] + list(tail)
            origins = f.origins
            if origins:
                o = origins[0]
                s["well_name"] = str(getattr(o, "well_name", "") or "")
                s["uwi"]       = str(getattr(o, "api_well", "") or
                                     getattr(o, "uwi", "") or "")
                company        = str(getattr(o, "company", "") or "")
                field          = str(getattr(o, "field_name", "") or "")
            else:
                company = field = ""

            total_ch = sum(len(lf.channels) for lf in lfs)
            total_fr = sum(len(lf.frames)   for lf in lfs)

            s["description"] = (
                f"DLIS · {len(lfs)} logical file(s) · "
                f"{total_ch} channels · {total_fr} frame(s)"
            )
            s["key_fields"] = {
                "logical_files": len(lfs),
                "channels":      total_ch,
                "frames":        total_fr,
                "company":       company,
                "field":         field,
            }
            s["ppdm_hints"] = ["dbo.WELL_LOG_SAMPLE", "dbo.WELL"]

    except Exception as e:
        s["error"] = str(e)
    return s


# ══════════════════════════════════════════════════════════════════════════════
# LIS
# ══════════════════════════════════════════════════════════════════════════════
def _summarize_lis(file_path: str) -> dict:
    s = _base(file_path, "LIS")
    try:
        from dataview.file_catalog.lis_catalog import classify_lis
        cl = classify_lis(file_path)

        s["well_name"] = cl.get("well_name")
        s["uwi"]       = cl.get("uwi")

        s["description"] = cl.get("description", "LIS file")
        s["key_fields"]  = {
            "well_name":   cl.get("well_name"),
            "uwi":         cl.get("uwi"),
            "operator":    cl.get("operator"),
            "field":       cl.get("well_field"),
            "state":       cl.get("state"),
            "county":      cl.get("county"),
            "contractor":  cl.get("contractor"),
            "depth_start": cl.get("depth_start"),
            "depth_stop":  cl.get("depth_stop"),
            "curves":      cl.get("n_curves", 0),
            "curve_names": cl.get("curve_names", []),
            "frames":      cl.get("n_frames", 0),
            "confidence":  cl.get("confidence", 0.0),
            "via_dlisio":  cl.get("via_dlisio", False),
        }
        s["ppdm_hints"] = ["dbo.WELL_LOG_SAMPLE", "dbo.WELL"]

        if not s["uwi"]:
            s["warnings"].append("No UWI/API found in LIS header")
        if cl.get("n_curves", 0) == 0:
            s["warnings"].append("No curve mnemonics extracted — file may be non-standard")
        if cl.get("error"):
            s["warnings"].append(f"Extraction note: {cl['error']}")

    except Exception as e:
        s["error"] = str(e)
    return s
def _summarize_segy(file_path: str) -> dict:
    s = _base(file_path, "SEG-Y")
    try:
        # Use the hand-rolled header-only reader (reads the first 3600 bytes +
        # samples ~300 trace headers for geometry) instead of segyio's
        # f.attributes(INLINE_3D)[:], which slurps EVERY trace header and crawls
        # on large files (verified: header-only does a 145k-trace 3D in 0.16s).
        from dataview.file_catalog.segy_header import read_segy_header
        h = read_segy_header(file_path)
        if not h.get("ok"):
            s["error"] = h.get("error") or "SEG-Y header read failed"
            return s

        n_traces  = h.get("n_traces") or 0
        si        = h.get("sample_interval_us")
        n_samples = h.get("n_samples")
        dims      = (h.get("dims") or "").replace("?", "")
        if dims not in ("2D", "3D"):
            dims = "3D" if n_traces > 10000 else "2D"
        ilr = h.get("inline_range")
        xlr = h.get("crossline_range")
        geom = (f"IL {ilr[0]}-{ilr[1]} · XL {xlr[0]}-{xlr[1]}"
                if ilr and xlr else "geometry not available")

        s["description"] = (
            f"SEG-Y · {dims} · {n_traces:,} traces"
            f" · {si} µs sample interval · {n_samples} samples · {geom}"
        )
        s["key_fields"] = {
            "traces":          n_traces,
            "seis_set_type":   dims,
            "sample_interval": si,
            "samples":         n_samples,
            "format":          h.get("format_desc"),
            "survey_name":     h.get("survey_name") if h.get("survey_name")
                               else None,
            "inline_range":    ilr,
            "crossline_range": xlr,
        }
        s["ppdm_hints"] = ["dbo.SEIS_SET", "dbo.SEIS_TRACE"]
        if n_traces == 0:
            s["warnings"].append("Zero traces — file may be corrupt")

    except Exception as e:
        s["error"] = str(e)
    return s


# ══════════════════════════════════════════════════════════════════════════════
# PDF
# ══════════════════════════════════════════════════════════════════════════════
def _summarize_pdf(file_path: str) -> dict:
    s = _base(file_path, "PDF")
    try:
        import pdfplumber
        from dataview.file_catalog.pdf_survey_catalog import classify_pdf

        cl = classify_pdf(file_path)
        # Single owner of PDF→fields resolution handles the classify +
        # extended-classify + scout-grid-header dance (incl. grid/GID UWIs).
        # We keep the base classify_pdf call above for station_count/description
        # fields the resolver doesn't carry, then overlay the resolved identity.
        try:
            from dataview.file_catalog.pdf_survey_catalog import resolve_pdf_fields
            r = resolve_pdf_fields(file_path)
            if r.get("report_type") and r["report_type"] != "UNKNOWN":
                cl["report_type"] = r["report_type"]
            for k in ("uwi", "well_name", "operator"):
                if r.get(k):
                    cl[k] = r[k]
        except Exception:
            pass

        s["well_name"] = cl.get("well_name")
        s["uwi"]       = cl.get("uwi")

        with pdfplumber.open(file_path) as pdf:
            pages    = len(pdf.pages)
            has_tbl  = any(pdf.pages[i].extract_tables()
                           for i in range(min(3, pages)))

        s["description"] = (
            f"PDF · {pages} page(s) · {cl['report_type'].replace('_',' ').title()}"
            f" · {cl.get('station_count',0)} stations"
            f" · {int(cl.get('confidence',0)*100)}% confidence"
        )
        s["key_fields"] = {
            "pages":        pages,
            "report_type":  cl["report_type"],
            "station_count":cl.get("station_count", 0),
            "confidence":   cl.get("confidence", 0),
            "has_tables":   has_tbl,
            "operator":     cl.get("operator"),
            "survey_type":  cl.get("survey_type"),
            "uwi":          cl.get("uwi"),
            "well_name":    cl.get("well_name"),
        }
        if cl["report_type"] == "DIRECTIONAL_SURVEY":
            s["ppdm_hints"] = ["dbo.WELL_DIR_SURVEY",
                               "dbo.WELL_DIR_SRVY_STATION"]
        elif cl["report_type"] == "MUD_LOG":
            s["ppdm_hints"] = ["dbo.WELL_LOG_SAMPLE"]

    except Exception as e:
        s["error"] = str(e)
    return s


# ══════════════════════════════════════════════════════════════════════════════
# Shapefile
# ══════════════════════════════════════════════════════════════════════════════
def _summarize_shp(file_path: str) -> dict:
    s = _base(file_path, "Shapefile")
    try:
        import geopandas as gpd
        from dataview.mapping.shapefile_catalog import classify_shapefile

        cl  = classify_shapefile(file_path)
        gdf = gpd.read_file(file_path, rows=10)

        n       = cl.get("feature_count", 0)
        geom    = cl.get("geometry_type", "?")
        crs     = cl.get("crs_epsg", "?")
        cols    = cl.get("attributes", [])
        ft      = cl.get("feature_type", "?")
        bounds  = cl.get("bounds") or {}
        col_map = cl.get("column_map", {})

        # ── Extract sample values from matched DBF columns ────────────────────
        sample_values: dict = {}

        # UWIs / well IDs
        if "UWI" in col_map:
            uwi_col = col_map["UWI"]
            if uwi_col in gdf.columns:
                vals = gdf[uwi_col].dropna().astype(str).tolist()
                sample_values["sample_uwis"] = [v for v in vals if v.strip()][:5]
                s["uwi"] = sample_values["sample_uwis"][0] \
                    if sample_values["sample_uwis"] else None

        # Well names
        if "WELL_NAME" in col_map:
            wn_col = col_map["WELL_NAME"]
            if wn_col in gdf.columns:
                vals = gdf[wn_col].dropna().astype(str).tolist()
                sample_values["sample_well_names"] = [v for v in vals if v.strip()][:5]
                s["well_name"] = sample_values["sample_well_names"][0] \
                    if sample_values["sample_well_names"] else None

        # Operators — load more rows to get a representative list
        if "OPERATOR" in col_map:
            op_col = col_map["OPERATOR"]
            try:
                full = gpd.read_file(file_path,
                                     include_fields=[op_col])
                ops = (full[op_col].dropna()
                                    .astype(str)
                                    .str.strip()
                                    .replace("", None)
                                    .dropna()
                                    .value_counts()
                                    .head(5)
                                    .index.tolist())
                sample_values["top_operators"] = ops
            except Exception:
                pass

        # Field names
        if "FIELD_NAME" in col_map:
            fn_col = col_map["FIELD_NAME"]
            if fn_col in gdf.columns:
                vals = gdf[fn_col].dropna().astype(str).tolist()
                sample_values["sample_fields"] = [v for v in vals if v.strip()][:5]

        # Date columns — extract range
        for date_key in ("SPUD_DATE", "COMPLETION_DATE"):
            if date_key in col_map:
                d_col = col_map[date_key]
                if d_col in gdf.columns:
                    try:
                        import pandas as pd
                        full_d = gpd.read_file(file_path,
                                               include_fields=[d_col])
                        dates = pd.to_datetime(
                            full_d[d_col], errors="coerce").dropna()
                        if len(dates):
                            sample_values[f"{date_key.lower()}_range"] = (
                                f"{dates.min().strftime('%Y-%m-%d')} – "
                                f"{dates.max().strftime('%Y-%m-%d')}"
                            )
                    except Exception:
                        pass

        s["description"] = (
            f"Shapefile · {n:,} {geom} features · {ft.replace('_',' ')}"
            f" · CRS EPSG:{crs} · {len(cols)} attributes"
        )
        s["key_fields"] = {
            "feature_count": n,
            "geometry_type": geom,
            "feature_type":  ft,
            "crs_epsg":      crs,
            "attributes":    cols[:10],
            "column_map":    col_map,
            "bounds":        bounds,
            "confidence":    cl.get("confidence", 0),
            **sample_values,
        }
        s["ppdm_hints"] = [cl.get("ppdm_target")] if cl.get("ppdm_target") else []

    except Exception as e:
        s["error"] = str(e)
    return s


# ══════════════════════════════════════════════════════════════════════════════
# Excel
# ══════════════════════════════════════════════════════════════════════════════

# Column pattern → PPDM table hints
EXCEL_TABLE_TYPES = {
    "PRODUCTION": {
        "keywords": ["oil","gas","water","bbls","mcf","boe",
                     "production","prod","volume","gross","net"],
        "ppdm":     "dbo.WELL_VERSION",
        "required": ["date","uwi"],
    },
    "COMPLETION": {
        "keywords": ["frac","stage","proppant","fluid","perf",
                     "cluster","completion","sand","water_vol"],
        "ppdm":     "dbo.WELL_COMPLETION",
        "required": ["uwi"],
    },
    "FORMATION_TOPS": {
        "keywords": ["formation","top","base","net_pay","pay",
                     "tops","horizon","marker","depth"],
        "ppdm":     "dbo.WELL_FORMATION",
        "required": ["uwi","formation"],
    },
    "WELL_HEADER": {
        "keywords": ["uwi","api","well_name","operator","latitude",
                     "longitude","field","county","kb","td"],
        "ppdm":     "dbo.WELL",
        "required": ["uwi"],
    },
    "CORE_ANALYSIS": {
        "keywords": ["porosity","permeability","perm","poro","sw",
                     "saturation","core","plug","grain"],
        "ppdm":     "dbo.WELL_CORE_ANALYSIS",
        "required": ["depth"],
    },
    "PRESSURE": {
        "keywords": ["pressure","psi","bhp","whp","temperature",
                     "gradient","datum","shut_in"],
        "ppdm":     "dbo.WELL_PRESSURE_SURVEY",
        "required": ["depth","pressure"],
    },
    "SURVEY": {
        "keywords": ["md","inc","azi","tvd","inclination","azimuth",
                     "measured_depth","dogleg","dls"],
        "ppdm":     "dbo.WELL_DIR_SRVY_STATION",
        "required": ["md","inc"],
    },
    "RESERVE": {
        "keywords": ["reserve","proved","probable","possible",
                     "1p","2p","3p","pdp","pud","mstb","mmcf"],
        "ppdm":     "dbo.WELL_VERSION",
        "required": ["uwi"],
    },
}


def _classify_excel_sheet(headers: list[str]) -> tuple[str, float]:
    """Classify a sheet by its column headers."""
    hdrs_lower = [h.lower().replace(' ','_').replace('/','_')
                  for h in headers if h]
    best_type  = "UNKNOWN"
    best_score = 0.0

    for table_type, cfg in EXCEL_TABLE_TYPES.items():
        score = sum(
            1 for kw in cfg["keywords"]
            if any(kw in h for h in hdrs_lower)
        )
        req_score = sum(
            1 for req in cfg["required"]
            if any(req in h for h in hdrs_lower)
        )
        # Weight required columns more heavily
        total = score + req_score * 2
        norm  = total / (len(cfg["keywords"]) + len(cfg["required"]) * 2)
        if norm > best_score:
            best_score = norm
            best_type  = table_type

    return best_type, round(best_score, 2)


# Known fixed-schema file patterns — checked before the generic column
# classifier so files with well-known structures get a precise classification.
KNOWN_SCHEMAS = {
    "BOEM_BOREHOLE": {
        "required": ["api well number", "well name", "bottom lease number",
                     "water depth"],
        "ppdm": "dbo.WELL",
        "description": "BOEM Gulf of Mexico Borehole Data",
    },
    "BOEM_WELL": {
        "required": ["api well number", "spud date", "status code",
                     "surface latitude"],
        "ppdm": "dbo.WELL",
        "description": "BOEM Well Header Data",
    },
    "KGS_WELL": {
        "required": ["api_number", "lease_name", "township", "range"],
        "ppdm": "dbo.WELL",
        "description": "Kansas Geological Survey Well Data",
    },
    "RRC_WELL": {
        "required": ["api14", "operator_name", "county_name", "district"],
        "ppdm": "dbo.WELL",
        "description": "Texas RRC Well Data",
    },
}


def _detect_known_schema(headers_lower: list) -> tuple | None:
    """Match column headers against known fixed schemas.

    Returns (schema_name, ppdm_target) when ALL required columns for a
    schema are found (case-insensitive substring match), else None.
    """
    for schema_name, cfg in KNOWN_SCHEMAS.items():
        if all(any(req in h for h in headers_lower)
               for req in cfg["required"]):
            return schema_name, cfg["ppdm"]
    return None


# ── Excel header-location heuristics ─────────────────────────────────────────
# The header row isn't always row 0: real exports carry a title, blank lines or
# metadata above it. We read a window of rows and pick the one that looks most
# like a header; if none does, we fall back to label:value forms.
_SCAN_ROWS = 25
# Tokens that, when several appear in one row, mark it as a header row.
_HEADER_HINTS = (
    "uwi", "api", "well", "operator", "field", "county", "state", "province",
    "country", "lat", "lon", "long", "formation", "top", "base", "depth",
    "date", "spud", "elevation", "name", "number", "id", "status", "type",
    "md", "tvd",
)
# Exact cell/label names that identify the well UWI/API value.
_UWI_LABELS = (
    "uwi", "uwi14", "api", "api14", "api number", "api_no", "apino",
    "api well number", "well_id", "well id", "wellid",
)
# Exact cell/label names that identify the well name value.
_WELLNAME_LABELS = ("well name", "well_name", "wellname", "well")


def _summarize_excel(file_path: str) -> dict:
    s = _base(file_path, "Excel")
    try:
        import openpyxl
        import pandas as pd

        # ── Single-pass streaming read ─────────────────────────────────────
        # openpyxl read_only + iter_rows() streams the XML row by row and
        # stops as soon as we break. This means a 50,000-row or 4-million-
        # row file reads exactly as fast as a 10-row file for our purposes —
        # we only need the header row and a handful of data rows.
        #
        # The previous approach called pd.read_excel(..., nrows=5) per sheet,
        # which re-opens and re-parses the entire worksheet XML on every call
        # regardless of nrows. On dense files this caused multi-minute hangs.
        _MAX_SAMPLE_ROWS = 5   # data rows to read for UWI / value sampling

        wb = openpyxl.load_workbook(
            file_path, read_only=True, data_only=True)
        sheets = wb.sheetnames

        sheet_summaries = []
        all_ppdm        = []
        total_rows      = 0
        uwi_found       = None
        well_name_found = None

        for sheet_name in sheets:
            try:
                ws = wb[sheet_name]

                # ── Locate the header row within the first _SCAN_ROWS rows ──
                # Read a window, then pick the row that looks most like a header
                # (most header-hint cells, at least two — so a one-word title
                # line isn't mistaken for a header).
                scan = []
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    scan.append(row)
                    if i + 1 >= _SCAN_ROWS:
                        break
                if not scan:
                    continue

                hdr_idx, best = None, 1
                for ri, row in enumerate(scan):
                    score = sum(
                        1 for c in row
                        if c is not None and str(c).strip()
                        and any(h in str(c).lower() for h in _HEADER_HINTS))
                    if score > best:
                        best, hdr_idx = score, ri

                if hdr_idx is not None:
                    headers     = [str(c) if c is not None else ""
                                   for c in scan[hdr_idx]]
                    sample_rows = scan[hdr_idx + 1:hdr_idx + 1 + _MAX_SAMPLE_ROWS]
                else:
                    headers     = [str(c) if c is not None else ""
                                   for c in scan[0]]
                    sample_rows = scan[1:1 + _MAX_SAMPLE_ROWS]
                headers_lower = [h.lower().strip() for h in headers]

                # Row count — subtract the header position so a title block
                # above the header doesn't inflate the count.
                _hpos      = (hdr_idx if hdr_idx is not None else 0) + 1
                n_rows     = max(0, (ws.max_row or _hpos) - _hpos)
                total_rows += n_rows

                # ── Schema detection: known fixed schemas first ─────────────
                known = _detect_known_schema(headers_lower)
                if known:
                    schema_name, ppdm = known
                    table_type = schema_name
                    conf       = 1.0
                else:
                    table_type, conf = _classify_excel_sheet(headers)
                    ppdm = EXCEL_TABLE_TYPES.get(
                        table_type, {}).get("ppdm", "")

                # ── UWI: a header column first, then a label:value form ─────
                if not uwi_found:
                    for col_i, col in enumerate(headers):
                        if any(x in col.lower() for x in
                               ("uwi", "api", "well_id", "well id",
                                "api well number")):
                            for row in sample_rows:
                                if col_i < len(row) and row[col_i] not in (None, ""):
                                    uwi_found = str(row[col_i]).strip()
                                    break
                            if uwi_found:
                                break
                if not uwi_found:
                    # label:value layout — a cell named UWI/API with its value in
                    # the next column. Runs even when a (sub-table) header was
                    # detected elsewhere on the sheet, but skips that header row
                    # and only accepts a UWI-like value (>=10 digits) so a header
                    # cell can never be mistaken for the value.
                    for ri, row in enumerate(scan):
                        if ri == hdr_idx:
                            continue
                        cells = list(row)
                        for ci in range(len(cells) - 1):
                            lab = (str(cells[ci]).lower().strip()
                                   if cells[ci] is not None else "")
                            if lab in _UWI_LABELS and cells[ci + 1] not in (None, ""):
                                cand = str(cells[ci + 1]).strip()
                                if sum(ch.isdigit() for ch in cand) >= 10:
                                    uwi_found = cand
                                    break
                        if uwi_found:
                            break

                # ── Well name: a header column first, then label:value ──────
                if not well_name_found:
                    for col_i, col in enumerate(headers):
                        if col.lower().strip() in _WELLNAME_LABELS:
                            for row in sample_rows:
                                if col_i < len(row) and row[col_i] not in (None, ""):
                                    well_name_found = str(row[col_i]).strip()[:120]
                                    break
                            if well_name_found:
                                break
                if not well_name_found:
                    for ri, row in enumerate(scan):
                        if ri == hdr_idx:
                            continue
                        cells = list(row)
                        for ci in range(len(cells) - 1):
                            lab = (str(cells[ci]).lower().strip()
                                   if cells[ci] is not None else "")
                            if lab in _WELLNAME_LABELS and cells[ci + 1] not in (None, ""):
                                val = str(cells[ci + 1]).strip()
                                if val.lower() not in _WELLNAME_LABELS + _UWI_LABELS:
                                    well_name_found = val[:120]
                                    break
                        if well_name_found:
                            break

                # ── Date range — flag presence only, no full-column read ───
                date_range = ""
                for col in headers:
                    if any(x in col.lower()
                           for x in ["date", "month", "year"]):
                        date_range = "date column present"
                        break

                if ppdm:
                    all_ppdm.append(ppdm)

                sheet_summaries.append({
                    "sheet":      sheet_name,
                    "table_type": table_type,
                    "confidence": conf,
                    "rows":       n_rows,
                    "columns":    len(headers),
                    "headers":    headers[:8],
                    "ppdm":       ppdm,
                    "date_range": date_range,
                })

            except Exception:
                pass

        wb.close()

        s["uwi"] = uwi_found
        if well_name_found:
            s["well_name"] = well_name_found
        s["description"] = (
            f"Excel · {len(sheets)} sheet(s) · {total_rows:,} total rows · "
            + " | ".join(
                f"{ss['sheet']}: {ss['table_type']} ({ss['rows']:,} rows)"
                for ss in sheet_summaries[:4]
            )
        )
        s["key_fields"] = {
            "sheets":       sheets,
            "total_rows":   total_rows,
            "sheet_detail": sheet_summaries,
        }
        s["ppdm_hints"] = list(dict.fromkeys(all_ppdm))

        if total_rows == 0:
            s["warnings"].append("No data rows found — file may be empty")
        if not s["ppdm_hints"]:
            s["warnings"].append("Could not classify sheet content")

    except Exception as e:
        s["error"] = str(e)
    return s


# ── Word document type taxonomy ──────────────────────────────────────────────
# _summarize_docx scores a document against these keyword sets to label its
# doc_type. Matching is substring (kw.upper() in text_up), so keep keywords
# clean and unambiguous. Order is irrelevant; the highest keyword-hit count wins.
WORD_DOC_TYPES = {
    "COMPLETION_REPORT":  ["completion", "perforation", "perforated", "frac",
                           "stimulation", "treatment", "well completion"],
    "DRILLING_REPORT":    ["drilling", "spud", "casing", "bit record",
                           "mud program", "drilling report", "tour"],
    "WELL_SUMMARY":       ["well summary", "summary report", "well report",
                           "well history", "wellbore summary"],
    "GEOLOGICAL_REPORT":  ["formation", "lithology", "geologic", "stratigraph",
                           "core analysis", "show report", "mud log"],
    "LOG_ANALYSIS":       ["log analysis", "petrophysical", "porosity",
                           "saturation", "net pay", "log interpretation"],
    "PRODUCTION_REPORT":  ["production report", "bopd", "mcf", "gor",
                           "decline", "allowable", "production history"],
    "WORKOVER_REPORT":    ["workover", "recompletion", "remedial",
                           "plug and abandon", "p&a", "fishing"],
    "SCOUT_TICKET":       ["scout ticket", "scout", "initial potential",
                           "ip test", "drillstem test"],
    "AFE":                ["authorization for expenditure", "afe",
                           "cost estimate", "dry hole cost"],
    "PERMIT":             ["permit", "application to drill", "w-1", "form 1",
                           "spacing", "regulatory"],
}


def _summarize_docx(file_path: str) -> dict:
    s = _base(file_path, "Word")
    try:
        import docx

        doc     = docx.Document(file_path)
        core    = doc.core_properties

        # Document properties
        title   = core.title   or ""
        author  = core.author  or ""
        created = str(core.created)[:10] if core.created else ""
        modified= str(core.modified)[:10] if core.modified else ""

        # Extract all text
        full_text = "\n".join(p.text for p in doc.paragraphs)
        text_up   = full_text.upper()

        # Extract headings
        headings = [p.text.strip() for p in doc.paragraphs
                    if p.style.name.startswith("Heading")
                    and p.text.strip()][:10]

        # Count tables and extract headers
        tables_info = []
        for i, tbl in enumerate(doc.tables[:10]):
            if not tbl.rows:
                continue
            hdr_row  = [c.text.strip() for c in tbl.rows[0].cells]
            n_rows   = len(tbl.rows) - 1
            tbl_type, conf = _classify_excel_sheet(hdr_row)
            tables_info.append({
                "table_idx":  i,
                "headers":    hdr_row[:8],
                "rows":       n_rows,
                "table_type": tbl_type,
                "confidence": conf,
                "ppdm":       EXCEL_TABLE_TYPES.get(tbl_type,{}).get("ppdm",""),
            })

        # Detect document type
        doc_type = "UNKNOWN"
        best_score = 0
        for dtype, keywords in WORD_DOC_TYPES.items():
            score = sum(1 for kw in keywords if kw.upper() in text_up)
            if score > best_score:
                best_score = score
                doc_type   = dtype

        # Extract UWI/well name from paragraph text first…
        uwi_match = re.search(
            r'(?:UWI|API)[:\s]+([0-9\-]{10,20})', full_text, re.IGNORECASE)
        well_match = re.search(
            r'(?:WELL\s+NAME|WELL)[:\s]+([A-Za-z0-9 #\-]+)',
            full_text, re.IGNORECASE)

        s["uwi"]       = uwi_match.group(1).strip() if uwi_match else None
        s["well_name"] = well_match.group(1).strip()[:50] if well_match else None

        # …then scan two-column Field/Value tables, which carry the
        # authoritative identity in well reports (e.g. a 'Well Identification'
        # table with rows like 'Well Name | ANADARKO 1H', 'UWI | 42-...').
        # These OVERRIDE the paragraph guesses, which often grab the document
        # TITLE ('WELL COMPLETION REPORT') as the well name.
        _kv = {}
        for tbl in doc.tables:
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells]
                if len(cells) == 2 and cells[0] and cells[1]:
                    _kv[cells[0].strip().lower()] = cells[1].strip()
        def _kvget(*labels):
            for lab in labels:
                if _kv.get(lab):
                    return _kv[lab]
            return None
        _t_uwi  = _kvget("uwi", "api", "api number", "api no", "well id")
        _t_name = _kvget("well name", "wellname", "well")
        if _t_uwi:
            s["uwi"] = _t_uwi
        if _t_name:
            s["well_name"] = _t_name
        # carry operator/field/state if the table has them (key_fields below)
        _t_operator = _kvget("operator", "company")
        _t_field    = _kvget("field", "field name")
        _t_state    = _kvget("state", "province", "province/state")

        ppdm_hints = list(dict.fromkeys(
            t["ppdm"] for t in tables_info if t.get("ppdm")
        ))

        s["description"] = (
            f"Word · {doc_type.replace('_',' ').title()}"
            f" · {len(doc.paragraphs)} paragraphs"
            f" · {len(doc.tables)} table(s)"
            f" · {len(full_text):,} characters"
        )
        s["key_fields"] = {
            "doc_type":     doc_type,
            "title":        title,
            "author":       author,
            "created":      created,
            "modified":     modified,
            "headings":     headings,
            "paragraphs":   len(doc.paragraphs),
            "tables":       len(doc.tables),
            "tables_detail":tables_info,
            "word_count":   len(full_text.split()),
            "operator":     _t_operator,
            "field":        _t_field,
            "state":        _t_state,
        }
        s["ppdm_hints"] = ppdm_hints

        if not doc.paragraphs:
            s["warnings"].append("Document appears empty")
        if best_score == 0:
            s["warnings"].append("Could not classify document type")

    except Exception as e:
        s["error"] = str(e)
    return s


# ══════════════════════════════════════════════════════════════════════════════
# CSV
# ══════════════════════════════════════════════════════════════════════════════
def _summarize_csv(file_path: str) -> dict:
    s = _base(file_path, "CSV")
    try:
        import pandas as pd

        # Read just headers + sample
        df   = pd.read_csv(file_path, nrows=5, low_memory=False)
        hdrs = [str(c) for c in df.columns]

        # Full row count without reading all data
        with open(file_path, 'r', errors='ignore') as f:
            n_rows = sum(1 for _ in f) - 1  # subtract header

        table_type, conf = _classify_excel_sheet(hdrs)
        ppdm = EXCEL_TABLE_TYPES.get(table_type, {}).get("ppdm","")

        # Find UWI
        for col in hdrs:
            if any(x in col.lower() for x in ['uwi','api','well_id']):
                if not df[col].empty:
                    s["uwi"] = str(df[col].iloc[0])
                break

        s["description"] = (
            f"CSV · {n_rows:,} rows · {len(hdrs)} columns · "
            f"{table_type.replace('_',' ').title()} ({int(conf*100)}%)"
        )
        s["key_fields"] = {
            "rows":       n_rows,
            "columns":    len(hdrs),
            "headers":    hdrs[:12],
            "table_type": table_type,
            "confidence": conf,
            "sample":     df.head(3).to_dict("records"),
        }
        s["ppdm_hints"] = [ppdm] if ppdm else []

        if n_rows == 0:
            s["warnings"].append("CSV has no data rows")

    except Exception as e:
        s["error"] = str(e)
    return s


# ══════════════════════════════════════════════════════════════════════════════
# P190
# ══════════════════════════════════════════════════════════════════════════════
def _summarize_p190(file_path: str) -> dict:
    s = _base(file_path, "P190")
    try:
        lines = []
        with open(file_path, 'r', errors='ignore') as f:
            lines = f.readlines()[:200]

        header_lines = [l for l in lines if l.startswith('H')]
        data_lines   = [l for l in lines
                        if l[0:1].upper() in ('S','T','C')]

        survey_name = ""
        vessel      = ""
        for hl in header_lines:
            if 'SURVEY' in hl.upper():
                survey_name = hl[2:].strip()[:60]
            if 'VESSEL' in hl.upper() or 'SHIP' in hl.upper():
                vessel = hl[2:].strip()[:40]

        # Full count
        with open(file_path, 'r', errors='ignore') as f:
            all_lines  = f.readlines()
        total_data = sum(1 for l in all_lines
                         if l[0:1].upper() in ('S','T','C'))

        s["description"] = (
            f"P190 · {len(header_lines)} header records"
            f" · {total_data:,} data records"
            f" · Survey: {survey_name or 'unknown'}"
        )
        s["key_fields"] = {
            "header_records": len(header_lines),
            "data_records":   total_data,
            "survey_name":    survey_name,
            "vessel":         vessel,
        }
        s["ppdm_hints"] = ["dbo.SEIS_LINE", "dbo.SEIS_SET"]

    except Exception as e:
        s["error"] = str(e)
    return s


# ══════════════════════════════════════════════════════════════════════════════
# ══════════════════════════════════════════════════════════════════════════════
# WITSML
# ══════════════════════════════════════════════════════════════════════════════
def _summarize_witsml(file_path: str) -> dict:
    s = _base(file_path, "WITSML")
    try:
        # Cheap namespace gate — don't parse non-WITSML XML
        with open(file_path, "rb") as fh:
            head = fh.read(500)
        if b"witsml.org/schemas" not in head:
            s["description"] = "XML file — not WITSML (no witsml.org namespace)"
            s["warnings"].append("Not a WITSML file — namespace not found in first 500 bytes")
            return s

        from dataview.file_catalog.witsml_catalog import classify_witsml
        cl = classify_witsml(file_path)

        s["well_name"] = cl.get("well_name")
        s["uwi"]       = cl.get("uwi")
        s["description"] = cl.get("description", "WITSML file")

        obj_type = cl.get("object_type", "unknown")
        s["key_fields"] = {
            "witsml_version":  cl.get("witsml_version"),
            "object_type":     obj_type,
            "n_objects":       cl.get("n_objects", 0),
            "well_name":       cl.get("well_name"),
            "uwi":             cl.get("uwi"),
            "operator":        cl.get("operator"),
            "contractor":      cl.get("contractor"),
            "depth_start":     cl.get("depth_start"),
            "depth_stop":      cl.get("depth_stop"),
            "confidence":      cl.get("confidence", 0.0),
        }

        # Object-type-specific fields
        if obj_type == "trajectory":
            s["key_fields"]["n_stations"] = cl.get("n_stations", 0)
            s["key_fields"]["survey_type"] = cl.get("survey_type")
            s["ppdm_hints"] = ["dbo.WELL_DIR_SURVEY", "dbo.WELL_DIR_SRVY_STATION"]
        elif obj_type == "log":
            s["key_fields"]["n_curves"]    = cl.get("n_curves", 0)
            s["key_fields"]["curve_names"] = cl.get("curve_names", [])
            s["ppdm_hints"] = ["dbo.WELL_LOG", "dbo.WELL_LOG_SAMPLE"]
        elif obj_type == "mudlog":
            s["key_fields"]["n_intervals"] = cl.get("n_intervals", 0)
            s["key_fields"]["gas_shows"]   = cl.get("gas_shows", [])
            s["ppdm_hints"] = ["dbo.WELL_LOG_SAMPLE"]
        else:
            s["ppdm_hints"] = ["dbo.WELL"]

        if cl.get("error"):
            s["warnings"].append(f"Extraction note: {cl['error']}")

    except Exception as e:
        s["error"] = str(e)
    return s


# ══════════════════════════════════════════════════════════════════════════════
# JSON Well Log / OSDU
# ══════════════════════════════════════════════════════════════════════════════
def _summarize_json_well_log(file_path: str) -> dict:
    s = _base(file_path, "JSON Well Log")
    try:
        # Cheap petroleum gate — skip non-petroleum JSON
        with open(file_path, "r", encoding="utf-8-sig", errors="replace") as fh:
            head_text = fh.read(512)
        _looks_petroleum = (
            '"kind"' in head_text or
            '"header"' in head_text or
            '"WellLog"' in head_text or
            '"wellbore"' in head_text.lower()
        )
        if not _looks_petroleum:
            s["description"] = "JSON file — not a recognised petroleum format"
            s["warnings"].append("No OSDU 'kind' or JSONWLF 'header' found in first 512 bytes")
            return s

        from dataview.file_catalog.json_well_log_catalog import classify_json_well_log
        cl = classify_json_well_log(file_path)

        s["well_name"] = cl.get("well_name")
        s["uwi"]       = cl.get("uwi")
        s["description"] = cl.get("description", "JSON petroleum file")

        schema = cl.get("json_schema", "unknown")
        s["key_fields"] = {
            "json_schema":  schema,
            "report_type":  cl.get("report_type"),
            "well_name":    cl.get("well_name"),
            "uwi":          cl.get("uwi"),
            "operator":     cl.get("operator"),
            "contractor":   cl.get("contractor"),
            "depth_start":  cl.get("depth_start"),
            "depth_stop":   cl.get("depth_stop"),
            "confidence":   cl.get("confidence", 0.0),
        }

        # Schema-specific additions
        if schema in ("osdu_well", "osdu_wellbore", "osdu_generic"):
            s["key_fields"]["latitude"]  = cl.get("latitude")
            s["key_fields"]["longitude"] = cl.get("longitude")
            s["key_fields"]["spud_date"] = cl.get("spud_date")
            s["key_fields"]["td_ft"]     = cl.get("total_depth")
            s["key_fields"]["county"]    = cl.get("county")
            s["key_fields"]["state"]     = cl.get("state")
            s["ppdm_hints"] = ["dbo.WELL"]
        elif schema == "osdu_well_log" or schema == "jsonwlf":
            s["key_fields"]["n_curves"]    = cl.get("n_curves", 0)
            s["key_fields"]["curve_names"] = cl.get("curve_names", [])
            s["ppdm_hints"] = ["dbo.WELL_LOG"]
        elif schema == "osdu_marker_set":
            s["key_fields"]["n_markers"]       = cl.get("n_markers", 0)
            s["key_fields"]["formation_names"] = cl.get("formation_names", [])
            s["key_fields"]["markers"]         = cl.get("markers", [])
            s["ppdm_hints"] = ["dbo.WELL_FORMATION_TOP"]
        elif schema == "osdu_pressure":
            s["key_fields"]["pressures"]      = cl.get("pressures", {})
            s["key_fields"]["n_flow_periods"] = cl.get("n_flow_periods", 0)
            s["key_fields"]["permeability"]   = cl.get("permeability")
            s["key_fields"]["fluid_type"]     = cl.get("fluid_type")
            s["ppdm_hints"] = ["dbo.WELL_TEST"]
        elif schema == "osdu_trajectory":
            s["key_fields"]["survey_params"] = cl.get("survey_params", {})
            s["key_fields"]["n_stations"]    = cl.get("n_stations", 0)
            s["ppdm_hints"] = ["dbo.WELL_DIR_SURVEY", "dbo.WELL_DIR_SRVY_STATION"]
        elif schema == "osdu_field":
            s["key_fields"]["field_params"] = cl.get("field_params", {})
            s["ppdm_hints"] = ["dbo.FIELD"]
        elif schema == "osdu_reservoir":
            s["key_fields"]["reservoir_params"] = cl.get("reservoir_params", {})
            s["ppdm_hints"] = ["dbo.RESERVOIR"]
        elif schema == "osdu_scal":
            s["key_fields"]["scal_params"] = cl.get("scal_params", {})
            s["ppdm_hints"] = ["dbo.WELL_CORE_ANALYSIS"]
        elif schema == "osdu_document":
            s["key_fields"]["doc_params"] = cl.get("doc_params", {})
            s["ppdm_hints"] = ["dbo.WELL"]
        elif schema == "osdu_horizon":
            s["key_fields"]["horizon_params"] = cl.get("horizon_params", {})
            s["ppdm_hints"] = ["dbo.SEIS_HORIZON"]
        elif schema == "osdu_fault":
            s["key_fields"]["fault_params"] = cl.get("fault_params", {})
            s["ppdm_hints"] = ["dbo.SEIS_FAULT"]
        elif schema == "osdu_completion":
            cp = cl.get("completion_params", {})
            s["key_fields"]["completion_params"] = cp
            s["ppdm_hints"] = ["dbo.WELL_COMPLETION"]
        elif schema == "osdu_core":
            s["key_fields"]["n_plugs"]    = cl.get("n_plugs", 0)
            s["key_fields"]["core_stats"] = cl.get("core_stats", {})
            s["key_fields"]["plugs"]      = cl.get("plugs", [])
            s["ppdm_hints"] = ["dbo.WELL_CORE"]
        elif schema == "osdu_production":
            s["key_fields"]["n_months"]            = cl.get("n_production_months", 0)
            s["key_fields"]["production_summary"]  = cl.get("production_summary", {})
            s["ppdm_hints"] = ["dbo.PDEN_WELL_PRODUCTION"]
        elif schema == "osdu_seismic":
            s["key_fields"]["survey_name"]  = cl.get("survey_name")
            s["key_fields"]["seis_set_type"]= cl.get("seis_set_type")
            s["key_fields"]["acq_params"]   = cl.get("acq_params", {})
            s["ppdm_hints"] = ["dbo.SEIS_SET"]
        else:
            s["ppdm_hints"] = ["dbo.WELL"]

        if cl.get("error"):
            s["warnings"].append(f"Extraction note: {cl['error']}")

    except Exception as e:
        s["error"] = str(e)
    return s


# Main dispatcher
# ══════════════════════════════════════════════════════════════════════════════
# ══════════════════════════════════════════════════════════════════════════════
# Added extractors (KML, PowerPoint, raster/image, vector, ASCII/deviation,
# email, OpenDocument, RTF). Each is self-contained and degrades to s["error"]
# on failure so promote never stalls.
# ══════════════════════════════════════════════════════════════════════════════
_TXT_UWI_RX  = re.compile(r'(?:UWI|API)[\s:#=]+([0-9][0-9\-]{8,19})', re.I)
_TXT_WELL_RX = re.compile(r'(?:WELL\s*NAME|WELL)[\s:=]+([A-Za-z0-9 #\-/\.]{2,60})', re.I)


def _text_identity(text: str):
    """Pull a (uwi, well_name) out of free text. UWI must carry >=8 digits so
    a stray word can't be mistaken for one."""
    uwi = wn = None
    if not text:
        return uwi, wn
    m = _TXT_UWI_RX.search(text)
    if m and sum(ch.isdigit() for ch in m.group(1)) >= 8:
        uwi = m.group(1).strip()
    m = _TXT_WELL_RX.search(text)
    if m:
        wn = m.group(1).strip().rstrip(" (-/")
        wn = re.split(r'\b(?:UWI|API)\b', wn, flags=re.I)[0].strip().rstrip(" (-/")
        if not wn or wn[:1].isdigit():     # that was probably the UWI, not a name
            wn = None
    return uwi, wn


def _summarize_kml(file_path: str) -> dict:
    s = _base(file_path, "KML")
    try:
        import zipfile, xml.etree.ElementTree as ET
        if file_path.lower().endswith(".kmz"):
            with zipfile.ZipFile(file_path) as z:
                name = next((n for n in z.namelist()
                             if n.lower().endswith(".kml")), None)
                raw = z.read(name) if name else b""
        else:
            with open(file_path, "rb") as f:
                raw = f.read()
        if not raw:
            s["error"] = "no KML content"
            return s
        txt = raw.decode("utf-8", "ignore")
        txt = re.sub(r'\sxmlns(:\w+)?="[^"]+"', '', txt)   # drop namespaces
        root = ET.fromstring(txt)
        placemarks = root.findall(".//Placemark")
        names = [(pm.findtext("name") or "").strip() for pm in placemarks]
        names = [n for n in names if n]
        ext = {}
        for d in root.findall(".//Data"):
            k = (d.get("name") or "").strip()
            v = (d.findtext("value") or "").strip()
            if k:
                ext[k] = v
        uwi = wn = None
        for k, v in ext.items():
            if k.lower() in ("uwi", "api", "well_id", "uwi14", "api14") \
                    and sum(c.isdigit() for c in v) >= 8:
                uwi = v
                break
        u2, w2 = _text_identity(" ".join(names) + " " + " ".join(ext.values()))
        uwi = uwi or u2
        wn = w2 or (names[0] if names else None)
        s["uwi"], s["well_name"] = uwi, wn
        s["key_fields"] = {"placemark_count": len(placemarks),
                           "sample_names": names[:5],
                           "ext_data_keys": list(ext.keys())[:10]}
        s["description"] = f"KML · {len(placemarks)} placemark(s)" + (
            f" · e.g. {names[0]}" if names else "")
    except Exception as e:
        s["error"] = str(e)
    return s


def _summarize_pptx(file_path: str) -> dict:
    s = _base(file_path, "PowerPoint")
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        chunks = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    chunks.append(shape.text_frame.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        chunks.append(" ".join(c.text for c in row.cells))
        full = "\n".join(c for c in chunks if c)
        s["uwi"], s["well_name"] = _text_identity(full)
        title = ""
        if len(prs.slides) and prs.slides[0].shapes.title is not None:
            title = prs.slides[0].shapes.title.text or ""
        s["key_fields"] = {"slide_count": len(prs.slides), "title": title[:120]}
        s["description"] = f"PowerPoint · {len(prs.slides)} slide(s)" + (
            f" · {title[:60]}" if title else "")
    except Exception as e:
        s["error"] = str(e)
    return s


def _summarize_raster(file_path: str) -> dict:
    ext = Path(file_path).suffix.lower()
    s = _base(file_path, "Raster")
    try:
        if ext in (".tif", ".tiff"):
            try:
                import rasterio
                with rasterio.open(file_path) as ds:
                    crs = ds.crs.to_epsg() if ds.crs else None
                    b = ds.bounds
                    s["format"] = "GeoTIFF" if crs else "TIFF"
                    s["key_fields"] = {
                        "width": ds.width, "height": ds.height, "bands": ds.count,
                        "crs_epsg": crs,
                        "bounds": ([round(b.left, 4), round(b.bottom, 4),
                                    round(b.right, 4), round(b.top, 4)]
                                   if crs else None)}
                    s["description"] = (
                        f"{'GeoTIFF' if crs else 'TIFF raster'} · "
                        f"{ds.width}\u00d7{ds.height} · {ds.count} band(s)"
                        + (f" · EPSG:{crs}" if crs else ""))
                    if not crs:
                        s["warnings"].append(
                            "TIFF — likely a scanned log/map; identity needs OCR")
                    return s
            except Exception:
                pass
        from PIL import Image
        with Image.open(file_path) as im:
            s["format"] = "Image"
            s["key_fields"] = {"width": im.width, "height": im.height,
                               "mode": im.mode, "codec": im.format}
            s["description"] = (f"Image · {im.width}\u00d7{im.height} · {im.mode}"
                                " · no text extraction (OCR not run)")
            s["warnings"].append("raster/image — identity needs OCR, not run here")
    except Exception as e:
        s["error"] = str(e)
    return s


def _summarize_vector(file_path: str) -> dict:
    """Generic GDAL/GeoPandas vector reader — MapInfo (.tab/.mif), file
    geodatabase (.gdb), and any other OGR vector format."""
    s = _base(file_path, "Vector")
    try:
        import geopandas as gpd
        gdf = gpd.read_file(file_path, rows=50)
        cols = list(gdf.columns)

        def find_col(keys):
            for c in cols:
                if str(c).lower() in keys or any(k in str(c).lower() for k in keys):
                    return c
            return None

        uc = find_col(("uwi", "api", "well_id", "wellid"))
        wc = find_col(("well_name", "wellname", "well_nm", "name"))
        if uc is not None:
            vals = [v for v in gdf[uc].dropna().astype(str)
                    if sum(c.isdigit() for c in v) >= 8]
            if vals:
                s["uwi"] = vals[0]
        if wc is not None:
            vals = [v for v in gdf[wc].dropna().astype(str) if v.strip()]
            if vals:
                s["well_name"] = vals[0]
        crs = gdf.crs.to_epsg() if gdf.crs is not None else None
        geom = gdf.geom_type.dropna().unique().tolist()
        try:
            total = len(gpd.read_file(file_path))
        except Exception:
            total = len(gdf)
        s["format"] = "Vector"
        s["key_fields"] = {"feature_count": total,
                           "geometry_type": geom[0] if geom else "?",
                           "crs_epsg": crs, "attributes": cols[:20]}
        s["description"] = (
            f"Vector · {total:,} {geom[0] if geom else '?'} feature(s) · "
            f"{len(cols)} attribute(s)" + (f" · EPSG:{crs}" if crs else ""))
    except Exception as e:
        s["error"] = str(e)
    return s


def _summarize_ascii_log(file_path: str) -> dict:
    ext = Path(file_path).suffix.lower()
    kind = "Deviation survey" if ext == ".dev" else "ASCII log"
    s = _base(file_path, kind)
    try:
        lines = []
        with open(file_path, "r", errors="ignore") as f:
            for i, ln in enumerate(f):
                lines.append(ln.rstrip("\n"))
                if i >= 400:
                    break
        head = "\n".join(lines[:80])
        uwi, wn = _text_identity(head)
        if not uwi:
            for ln in lines[:80]:
                m = re.match(r'\s*(?:UWI|API|WELL[_ ]?ID)\b[\s:=]+'
                             r'([0-9][0-9\-]{8,19})', ln, re.I)
                if m and sum(c.isdigit() for c in m.group(1)) >= 8:
                    uwi = m.group(1).strip()
                    break
        if not wn:
            for ln in lines[:80]:
                m = re.match(r'\s*WELL(?:[_ ]?NAME)?\b[\s:=]+(.+)', ln, re.I)
                if m and m.group(1).strip()[:1] and not m.group(1).strip()[0].isdigit():
                    wn = m.group(1).strip()[:60]
                    break
        ncols = 0
        for ln in lines:
            toks = ln.split()
            nums = sum(1 for t in toks if re.fullmatch(r'[-+]?\d+(?:\.\d+)?', t))
            if nums >= 3:
                ncols = max(ncols, len(toks))
        s["uwi"], s["well_name"] = uwi, wn
        s["key_fields"] = {"line_count": len(lines), "data_columns": ncols}
        s["description"] = f"{kind} · ~{len(lines)} lines · {ncols} data column(s)"
    except Exception as e:
        s["error"] = str(e)
    return s


def _summarize_email(file_path: str) -> dict:
    ext = Path(file_path).suffix.lower()
    s = _base(file_path, "Email")
    try:
        subject = body = sender = ""
        if ext == ".eml":
            import email
            from email import policy
            with open(file_path, "rb") as f:
                msg = email.message_from_binary_file(f, policy=policy.default)
            subject = msg.get("subject", "") or ""
            sender = msg.get("from", "") or ""
            part = msg.get_body(preferencelist=("plain", "html"))
            body = part.get_content() if part else ""
        else:                                   # .msg
            import extract_msg
            m = extract_msg.Message(file_path)
            subject = m.subject or ""
            sender = m.sender or ""
            body = m.body or ""
        s["uwi"], s["well_name"] = _text_identity(f"{subject}\n{body}")
        s["key_fields"] = {"subject": subject[:120], "from": sender[:120]}
        s["description"] = (f"Email · {subject[:80]}" if subject else "Email")
    except Exception as e:
        s["error"] = str(e)
    return s


def _summarize_odf(file_path: str) -> dict:
    s = _base(file_path, "OpenDocument")
    try:
        import zipfile
        with zipfile.ZipFile(file_path) as z:
            xml = z.read("content.xml").decode("utf-8", "ignore")
        text = re.sub(r"<[^>]+>", " ", xml)
        text = re.sub(r"\s+", " ", text).strip()
        s["uwi"], s["well_name"] = _text_identity(text)
        s["key_fields"] = {"char_count": len(text)}
        s["description"] = f"OpenDocument · {len(text):,} chars"
    except Exception as e:
        s["error"] = str(e)
    return s


def _summarize_rtf(file_path: str) -> dict:
    s = _base(file_path, "RTF")
    try:
        from striprtf.striprtf import rtf_to_text
        with open(file_path, "r", errors="ignore") as f:
            text = rtf_to_text(f.read())
        s["uwi"], s["well_name"] = _text_identity(text)
        s["key_fields"] = {"char_count": len(text)}
        s["description"] = f"RTF · {len(text):,} chars"
    except Exception as e:
        s["error"] = str(e)
    return s


_DISPATCH = {
    ".las":    _summarize_las,
    ".dlis":   _summarize_dlis,
    ".dlf":    _summarize_dlis,
    ".lis":    _summarize_lis,
    ".segy":   _summarize_segy,
    ".sgy":    _summarize_segy,
    ".seg":    _summarize_segy,
    ".pdf":    _summarize_pdf,
    ".shp":    _summarize_shp,
    ".geojson":_summarize_shp,
    ".gpkg":   _summarize_shp,
    ".xml":    _summarize_witsml,
    ".json":   _summarize_json_well_log,
    ".xlsx":   _summarize_excel,
    ".xls":    _summarize_excel,
    ".docx":   _summarize_docx,
    ".doc":    _summarize_docx,
    ".csv":    _summarize_csv,
    ".txt":    _summarize_csv,
    ".tsv":    _summarize_csv,
    ".p190":   _summarize_p190,
    ".p90":    _summarize_p190,
    ".p1":     _summarize_p190,
    ".p2":     _summarize_p190,
    ".p3":     _summarize_p190,
    ".wml":    _summarize_witsml,
    ".xlsm":   _summarize_excel,
    ".docm":   _summarize_docx,
    ".kml":    _summarize_kml,
    ".kmz":    _summarize_kml,
    ".pptx":   _summarize_pptx,
    ".ppt":    _summarize_pptx,
    ".tif":    _summarize_raster,
    ".tiff":   _summarize_raster,
    ".png":    _summarize_raster,
    ".jpg":    _summarize_raster,
    ".jpeg":   _summarize_raster,
    ".tab":    _summarize_vector,
    ".mif":    _summarize_vector,
    ".gdb":    _summarize_vector,
    ".asc":    _summarize_ascii_log,
    ".prn":    _summarize_ascii_log,
    ".dev":    _summarize_ascii_log,
    ".eml":    _summarize_email,
    ".msg":    _summarize_email,
    ".odt":    _summarize_odf,
    ".ods":    _summarize_odf,
    ".odp":    _summarize_odf,
    ".rtf":    _summarize_rtf,
}

# The set of extensions this module can extract from. The workbench imports
# this to decide whether to route a file through summarize().
SUPPORTED_EXTS = frozenset(_DISPATCH)


def summarize(file_path: str) -> dict:
    """
    Universal summarizer — dispatches to the correct format handler.
    Always returns a dict with at minimum:
        file_path, file_name, format, size_kb,
        well_name, uwi, description, key_fields,
        warnings, ppdm_hints, error
    """
    ext = Path(file_path).suffix.lower()

    handler = _DISPATCH.get(ext)
    if handler:
        return handler(file_path)

    # Unknown format
    s = _base(file_path, ext.lstrip(".").upper() or "UNKNOWN")
    s["description"] = f"Unsupported format: {ext}"
    s["warnings"].append(f"No summarizer available for {ext} files")
    return s
